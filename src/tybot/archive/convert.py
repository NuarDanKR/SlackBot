"""첨부 문서 → 텍스트 변환.

실사용자가 올리는 건 xlsx·pdf·한글·워드·PPT 다. 개발자 텍스트 파일만 수집하면
맥락의 대부분을 놓친다. 그래서 변환한다 — 단 아래 규칙을 지킨다.

## 규칙
1. **변환본은 원문 라인으로 들어가되 `[첨부추출:파일명]` 로 표시한다.**
   사람이 "이건 자동 변환된 텍스트"임을 항상 알 수 있어야 한다.
2. **표는 구조를 살려 옮긴다**(시트명·행 단위). 요약·정리하지 않는다 — 그건 LLM 이 할 일이다.
3. **변환 실패는 조용히 넘기지 않는다.** 목록 줄은 남기고 경고를 올린다.
4. **스캔 PDF(텍스트 레이어 없음)는 변환하지 않는다.** OCR 은 오류가 사실처럼 굳는 경로다.
5. 라이브러리가 없으면 미변환으로 처리한다. 설치 여부가 조용한 고장이 되지 않게 경고를 남긴다.

## 지원
| 형식 | 방법 | 비고 |
|---|---|---|
| xlsx/xlsm | openpyxl (values only) | 수식 결과값. 시트·행 단위 |
| docx | python-docx | 문단 + 표 |
| pptx | python-pptx | 슬라이드별 텍스트 프레임 |
| pdf | pypdf | 텍스트 레이어만. 스캔본은 미변환 |
| hwpx | zipfile + XML | 한글 2014+ 표준 포맷 |
| hwp(구형 바이너리) | 미변환 | 신뢰할 만한 순수 파이썬 파서가 없다 |
| 이미지·도면 | 미변환 | OCR 미도입 |

## XML 안전
hwpx 는 사용자가 올린 zip 안의 XML 이다. 표준 파서는 외부 엔티티(XXE)와
엔티티 폭탄(billion laughs)에 취약하므로 `defusedxml` 로만 파싱한다.
없으면 변환하지 않는다(fail closed).
"""
from __future__ import annotations

import io
import logging
import re
import zipfile

logger = logging.getLogger("tybot.convert")

MAX_LINES = 400  # 파일 하나가 아카이브를 잡아먹지 않게
MAX_CELL = 200  # 셀 한 칸 길이 상한
CONVERTIBLE = {"xlsx", "xlsm", "docx", "pptx", "pdf", "hwpx"}


class ConvertError(RuntimeError):
    """변환 실패. 목록 줄은 남기고 경고로 올린다."""


def _clip(s: object) -> str:
    t = str(s).replace("\r", " ").replace("\n", " ").strip()
    return t if len(t) <= MAX_CELL else t[:MAX_CELL] + "…"


def _finish(lines: list[str]) -> list[str]:
    if len(lines) > MAX_LINES:
        lines = [*lines[:MAX_LINES], f"…(이하 생략, 총 {len(lines)}줄)"]
    return [ln for ln in lines if ln.strip()]


def _xlsx(data: bytes) -> list[str]:
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise ConvertError("openpyxl 미설치 - pip install openpyxl") from e

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    out: list[str] = []
    try:
        for ws in wb.worksheets:
            out.append(f"[시트] {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [_clip(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    out.append(" | ".join(cells))
                if len(out) > MAX_LINES + 5:
                    break
    finally:
        wb.close()
    return _finish(out)


def _docx(data: bytes) -> list[str]:
    try:
        import docx
    except ImportError as e:
        raise ConvertError("python-docx 미설치 - pip install python-docx") from e

    d = docx.Document(io.BytesIO(data))
    out = [_clip(p.text) for p in d.paragraphs if p.text.strip()]
    for ti, table in enumerate(d.tables, 1):
        out.append(f"[표 {ti}]")
        for row in table.rows:
            cells = [_clip(c.text) for c in row.cells if c.text.strip()]
            if cells:
                out.append(" | ".join(cells))
    return _finish(out)


def _pptx(data: bytes) -> list[str]:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ConvertError("python-pptx 미설치 - pip install python-pptx") from e

    prs = Presentation(io.BytesIO(data))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"[슬라이드 {i}]")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for p in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in p.runs).strip()
                    if txt:
                        out.append(_clip(txt))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [_clip(c.text) for c in row.cells if c.text.strip()]
                    if cells:
                        out.append(" | ".join(cells))
    return _finish(out)


def _pdf(data: bytes) -> list[str]:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise ConvertError("pypdf 미설치 - pip install pypdf") from e

    reader = PdfReader(io.BytesIO(data))
    if getattr(reader, "is_encrypted", False):
        raise ConvertError("암호가 걸린 PDF - 변환하지 않음")
    out: list[str] = []
    for i, page in enumerate(reader.pages, 1):
        try:
            text = page.extract_text() or ""
        except Exception:  # noqa: BLE001 - 한 페이지 실패가 전체를 막지 않는다
            continue
        lines = [_clip(ln) for ln in text.splitlines() if ln.strip()]
        if lines:
            out.append(f"[{i}쪽]")
            out.extend(lines)
    if not out:
        # 텍스트 레이어가 없다 = 스캔본. OCR 은 도입하지 않는다(원칙 4).
        raise ConvertError("텍스트 레이어 없음(스캔본으로 보임) - 원본 확인 필요")
    return _finish(out)


HWPX_TEXT_TAGS = ("t", "char")


def _hwpx(data: bytes) -> list[str]:
    """hwpx 는 zip + XML 이다. 텍스트 노드만 순서대로 뽑는다.

    파싱은 defusedxml 로만 한다 - 사용자가 올린 XML 이므로 XXE·엔티티 폭탄 대상이다.
    """
    try:
        from defusedxml.common import DefusedXmlException
        from defusedxml.ElementTree import ParseError, fromstring
    except ImportError as e:
        raise ConvertError(
            "defusedxml 미설치 - 사용자 XML 을 표준 파서로 열지 않는다. pip install defusedxml"
        ) from e

    out: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            names = [n for n in z.namelist() if n.startswith("Contents/") and n.endswith(".xml")]
            if not names:
                raise ConvertError("hwpx 구조가 예상과 다르다")
            for name in sorted(names):
                try:
                    root = fromstring(z.read(name))
                except ParseError:
                    continue
                except DefusedXmlException as e:
                    # 외부 엔티티·엔티티 폭탄 등. 방어했다는 사실을 남기고 파일 전체를 거부한다.
                    logger.warning("악성 XML 구조 감지 - 변환 거부: %s", e)
                    raise ConvertError(f"안전하지 않은 XML 구조: {e.__class__.__name__}") from e
                buf: list[str] = []
                for el in root.iter():
                    tag = el.tag.rsplit("}", 1)[-1]
                    if tag in HWPX_TEXT_TAGS and el.text and el.text.strip():
                        buf.append(el.text.strip())
                if buf:
                    joined = re.sub(r"\s{2,}", " ", " ".join(buf))
                    out.extend(_clip(ln) for ln in joined.split("\n"))
    except zipfile.BadZipFile as e:
        raise ConvertError("hwpx 파일이 손상됐거나 구형 hwp 형식") from e
    if not out:
        raise ConvertError("hwpx 에서 텍스트를 찾지 못했다")
    return _finish(out)


_HANDLERS = {
    "xlsx": _xlsx, "xlsm": _xlsx,
    "docx": _docx,
    "pptx": _pptx,
    "pdf": _pdf,
    "hwpx": _hwpx,
}


def can_convert(filetype: str) -> bool:
    return filetype.lower() in _HANDLERS


def convert(filetype: str, data: bytes) -> list[str]:
    """확장자별 변환. 실패는 ConvertError 로 올린다."""
    fn = _HANDLERS.get(filetype.lower())
    if fn is None:
        raise ConvertError(f"'{filetype}' 은 변환 대상이 아니다")
    if not data:
        raise ConvertError("빈 파일")
    return fn(data)
