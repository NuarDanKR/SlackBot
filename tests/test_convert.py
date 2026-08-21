"""문서 변환 — 실제 파일을 만들어 검증한다(모킹 아님)."""
from __future__ import annotations

import io
import zipfile

import pytest

from tybot.archive.convert import ConvertError, can_convert, convert


def test_xlsx_keeps_sheet_and_row_structure():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "기성내역"
    ws.append(["현장", "금액", "비고"])
    ws.append(["김해외동", "3억 2천만원", ""])
    ws2 = wb.create_sheet("요약")
    ws2.append(["합계", "3억 2천만원"])
    buf = io.BytesIO()
    wb.save(buf)

    lines = convert("xlsx", buf.getvalue())
    assert "[시트] 기성내역" in lines
    assert "현장 | 금액 | 비고" in lines
    assert "김해외동 | 3억 2천만원" in lines
    assert "[시트] 요약" in lines  # 시트가 여러 개면 모두 나온다


def test_docx_paragraphs_and_tables():
    import docx

    d = docx.Document()
    d.add_paragraph("착공일은 2026-03-01 이다")
    tbl = d.add_table(rows=2, cols=2)
    tbl.cell(0, 0).text = "항목"
    tbl.cell(0, 1).text = "값"
    tbl.cell(1, 0).text = "기성금"
    tbl.cell(1, 1).text = "3억"
    buf = io.BytesIO()
    d.save(buf)

    lines = convert("docx", buf.getvalue())
    assert "착공일은 2026-03-01 이다" in lines
    assert "[표 1]" in lines
    assert "기성금 | 3억" in lines


def test_pptx_slides():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "주간 보고"
    buf = io.BytesIO()
    prs.save(buf)

    lines = convert("pptx", buf.getvalue())
    assert "[슬라이드 1]" in lines
    assert "주간 보고" in lines


def test_pdf_without_text_layer_is_refused():
    """스캔본은 변환하지 않는다 - OCR 오류가 사실처럼 굳는 걸 막는다."""
    from pypdf import PdfWriter

    w = PdfWriter()
    w.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    w.write(buf)

    with pytest.raises(ConvertError, match="텍스트 레이어 없음"):
        convert("pdf", buf.getvalue())


def test_hwpx_text_extraction():
    xml = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<hml xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">'
        "<p><hp:t>착공 지연 사유</hp:t></p><p><hp:t>자재 수급 문제</hp:t></p></hml>"
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr("Contents/section0.xml", xml)
    lines = convert("hwpx", buf.getvalue())
    assert any("착공 지연 사유" in ln for ln in lines)
    assert any("자재 수급 문제" in ln for ln in lines)


def test_hwpx_billion_laughs_is_refused():
    """사용자가 올린 XML 이므로 엔티티 폭탄을 막아야 한다."""
    bomb = (
        '<?xml version="1.0"?><!DOCTYPE lolz [<!ENTITY lol "lol">'
        '<!ENTITY lol2 "&lol;&lol;&lol;&lol;&lol;&lol;&lol;&lol;&lol;&lol;">'
        ']><hml><p><t>&lol2;</t></p></hml>'
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr("Contents/section0.xml", bomb)
    # defusedxml 이 엔티티 선언 자체를 거부한다 → 텍스트를 못 찾고 실패
    with pytest.raises(ConvertError):
        convert("hwpx", buf.getvalue())


def test_old_hwp_binary_is_not_convertible():
    assert not can_convert("hwp")
    with pytest.raises(ConvertError, match="변환 대상이 아니다"):
        convert("hwp", b"\xd0\xcf\x11\xe0binary")


def test_empty_file_refused():
    with pytest.raises(ConvertError, match="빈 파일"):
        convert("xlsx", b"")


def test_convertible_set():
    for ext in ("xlsx", "xlsm", "docx", "pptx", "pdf", "hwpx"):
        assert can_convert(ext)
    for ext in ("hwp", "xls", "doc", "ppt", "jpg", "dwg", "zip"):
        assert not can_convert(ext)


def test_large_sheet_is_truncated():
    from openpyxl import Workbook

    wb = Workbook()
    for i in range(600):
        wb.active.append([f"행 {i}"])
    buf = io.BytesIO()
    wb.save(buf)
    lines = convert("xlsx", buf.getvalue())
    assert len(lines) <= 402
    assert any("이하 생략" in ln for ln in lines)
